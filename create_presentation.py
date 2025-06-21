from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from urllib.request import urlretrieve

# Create a presentation with a professional theme
prs = Presentation()

# Set up consistent styling
def create_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Set title styling
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
    
    # Set subtitle styling
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    return slide

def create_analysis_slide(prs, title, image_url, description):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Download and add image
    image_path = f"temp_{title.replace(' ', '_')}.png"
    urlretrieve(image_url, image_path)
    left = Inches(1)
    top = Inches(1.5)
    width = height = Inches(7)
    slide.shapes.add_picture(image_path, left, top, width, height)
    
    # Add description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(8), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    for paragraph in desc_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(12)
    
    return image_path

def create_text_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = content
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        paragraph.space_after = Pt(12)
    
    return slide

# Create title slide
title_slide = create_title_slide(
    prs,
    "Exploratory Data Analysis: Cars Dataset",
    "A comprehensive analysis of car specifications and pricing patterns"
)

# Add introduction slide
intro_slide = create_text_slide(
    prs,
    "Project Overview",
    "This presentation contains a comprehensive Exploratory Data Analysis (EDA) of a car dataset, examining various aspects of car specifications and pricing patterns."
)

# Add slides for each visualization
price_distribution_slide = create_analysis_slide(
    prs,
    "Price Distribution Analysis",
    "https://raw.githubusercontent.com/lakshya189/Exploratory-Data-Analysis-Cars-Dataset-Assignment/master/visualizations/price_distribution.png",
    "Shows the distribution of car prices across different market segments\nHighlights the concentration of prices in different ranges\nIdentifies premium and luxury vehicle segments"
)

hp_vs_price_slide = create_analysis_slide(
    prs,
    "Engine Power vs Price",
    "https://raw.githubusercontent.com/lakshya189/Exploratory-Data-Analysis-Cars-Dataset-Assignment/master/visualizations/hp_vs_price.png",
    "Demonstrates the relationship between engine power and vehicle price\nShows how higher HP generally correlates with higher prices\nIdentifies price thresholds at different HP levels"
)

mpg_relationship_slide = create_analysis_slide(
    prs,
    "Fuel Efficiency Analysis",
    "https://raw.githubusercontent.com/lakshya189/Exploratory-Data-Analysis-Cars-Dataset-Assignment/master/visualizations/mpg_relationship.png",
    "Shows the relationship between highway and city fuel efficiency\nHighlights the strong positive correlation\nIdentifies efficiency clusters across different vehicle types"
)

price_by_make_slide = create_analysis_slide(
    prs,
    "Market Segmentation",
    "https://raw.githubusercontent.com/lakshya189/Exploratory-Data-Analysis-Cars-Dataset-Assignment/master/visualizations/price_by_make.png",
    "Shows price distribution across different car makes\nHighlights distinct market segments\nDemonstrates pricing strategies across brands"
)

# Add conclusion slide
conclusion_slide = create_text_slide(
    prs,
    "Key Findings",
    "1. Price Analysis:\n   - Clear distribution of car prices\n   - Strong correlation between engine power and price\n   - Distinct market segments by make\n\n2. Engine Analysis:\n   - Higher horsepower generally leads to higher prices\n   - Clear price thresholds at different HP levels\n\n3. Fuel Efficiency:\n   - Strong positive correlation between city and highway MPG\n   - Efficiency clusters indicate different vehicle types"
)

findings_box = conclusion_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
findings_frame = findings_box.text_frame
findings_frame.text = "Main Findings:\n\n• Clear market segmentation based on\n  price and specifications\n• Strong correlation between technical\n  specifications and pricing\n• Distinct patterns in fuel efficiency\n  across different vehicle types\n• Significant variation in pricing\n  strategies across car makes"

# Add recommendations
rec_box = conclusion_slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(2.5))
rec_frame = rec_box.text_frame
rec_frame.text = "Recommendations:\n\n1. For Manufacturers:\n   • Optimize pricing strategy based on\n     market positioning\n   • Consider fuel efficiency in pricing\n     decisions\n\n2. For Consumers:\n   • Consider both price and\n     specifications when purchasing\n   • Understand market segments\n     and price ranges\n\n3. For Further Analysis:\n   • Explore additional features\n   • Analyze temporal trends\n   • Investigate regional pricing\n     variations"

# Save the presentation
prs.save('Car_Analysis_Presentation.pptx')
print("Professional presentation created successfully!")

# Add image slide for MPG relationship
slide = prs.slides.add_slide(prs.slide_layouts[5])
left = top = Inches(1)
width = height = Inches(7)
slide.shapes.add_picture('mpg_relationship.png', left, top, width, height)

txtbox = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(8), Inches(1))
txt_frame = txtbox.text_frame
txt_frame.text = "MPG Relationship:\n- Shows relationship between highway and city MPG\n- Strong positive correlation"

# Add image slide for Price by Make
slide = prs.slides.add_slide(prs.slide_layouts[5])
left = top = Inches(1)
width = height = Inches(7)
slide.shapes.add_picture('price_by_make.png', left, top, width, height)

txtbox = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(8), Inches(1))
txt_frame = txtbox.text_frame
txt_frame.text = "Price Distribution by Make:\n- Shows price variation across different car makes\n- Luxury brands have distinct price distributions"

# Add conclusions slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
txtbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
txt_frame = txtbox.text_frame
txt_frame.text = "Key Findings:\n\n1. Price Analysis:\n- Clear price distribution patterns\n- Higher HP engines correlate with higher prices\n- Luxury brands have distinct pricing\n\n2. Engine Analysis:\n- Wide range of engine specifications\n- 4-cylinder engines are most common\n\n3. Fuel Efficiency:\n- Strong correlation between highway and city MPG\n- Most cars have highway MPG between 20-30 mpg\n\n4. Market Segments:\n- Clear distinction between economy and luxury segments\n- Different pricing strategies across makes"

# Save the presentation
prs.save('Car_Analysis_Presentation.pptx')
print("Presentation created successfully!")
